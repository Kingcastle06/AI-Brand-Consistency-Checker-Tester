import streamlit as st
import os
import sys
import re
import io
import time
import uuid
from urllib.request import urlopen, Request

# --- DIRECT BACKEND INTEGRATION LAYER ---
# Adds the subfolder path to Python's memory search list so it can find orchestrator.py
backend_path = os.path.join(os.path.abspath(os.path.dirname(__file__)), "ai_brand_consistency_checker")
if os.path.exists(backend_path) and backend_path not in sys.path:
    sys.path.insert(0, backend_path)

try:
    import orchestrator as backend_engine
except ImportError:
    backend_engine = None

# --- NATIVE FILE PARSING UTILITIES ---
def extract_text_from_html(html_content):
    """Strips HTML tags to extract raw web page text natively."""
    clean_text = re.sub(r'<script.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
    clean_text = re.sub(r'<style.*?</style>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
    clean_text = re.sub(r'<[^>]+>', ' ', clean_text)
    clean_text = re.sub(r'\s+', ' ', clean_text).strip()
    return clean_text

def extract_text_from_pdf(file_bytes):
    """Extracts text strings out of binary PDF content safely."""
    try:
        import pypdf
        pdf_stream = io.BytesIO(file_bytes)
        reader = pypdf.PdfReader(pdf_stream)
        extracted_pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                extracted_pages.append(text)
        return "\n".join(extracted_pages)
    except Exception as e:
        return f"[PDF Parsing Fallback due to string structure exception: {str(e)}]"

def extract_text_from_docx(file_bytes):
    """Extracts paragraphs out of Microsoft Word documents safely."""
    try:
        import docx
        doc_stream = io.BytesIO(file_bytes)
        doc = docx.Document(doc_stream)
        return "\n".join([p.text for p in doc.paragraphs if p.text])
    except Exception as e:
        return f"[Word Doc Parsing Fallback due to format exception: {str(e)}]"

def extract_text_from_pptx(file_bytes):
    """Extracts structural text elements out of PowerPoint slides safely."""
    try:
        import pptx
        ppt_stream = io.BytesIO(file_bytes)
        prs = pptx.Presentation(ppt_stream)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        return f"[PowerPoint Parsing Fallback due to layout exception: {str(e)}]"

def highlight_violations(source_text, issues_list):
    """Highlights flagged brand violations directly in the source text copy."""
    html_markup = source_text
    for issue in issues_list:
        offending = issue.get("offending_text", "")
        if offending and offending.strip() and offending in html_markup:
            escaped_offending = re.escape(offending)
            highlighted_patch = f'<mark style="background-color: #FF007F; color: white; font-weight: bold; padding: 2px 4px; border-radius: 4px;">{offending}</mark>'
            try:
                html_markup = re.sub(escaped_offending, highlighted_patch, html_markup, flags=re.IGNORECASE)
            except Exception:
                continue
    return html_markup

# --- USER INTERFACE CANVAS CONFIGURATION ---
st.set_page_config(page_title="AI Brand Compliance Dashboard", layout="centered")

st.title("AI Brand Consistency Checker")
st.markdown("Upload a document/asset or provide a live website URL to run an automated compliance audit against your custom brand rules.")

st.sidebar.header("📋 Brand Guidelines")
rules_input = st.sidebar.text_area(
    "Enter Corporate Rules / Colors:",
    value="The only allowed colors are Cyber Pink #FF007F and Matrix Green #39FF14. We are testing neon styles."
)

input_mode = st.radio("Select Asset Target Type:", ["Uploaded Media File", "Live Website URL"], horizontal=True)

uploaded_file = None
url_input = None

if input_mode == "Uploaded Media File":
    uploaded_file = st.file_uploader(
        "Drag and drop your asset here...", 
        type=["pdf", "docx", "pptx", "png", "jpg", "jpeg", "txt"]
    )
else:
    url_input = st.text_input("Enter the corporate website URL to audit (e.g., https://example.com/landing):")

has_target = (uploaded_file is not None) if input_mode == "Uploaded Media File" else (url_input and url_input.strip() != "")

if has_target:
    if input_mode == "Uploaded Media File":
        st.info(f"📂 File loaded: {uploaded_file.name}")
    else:
        st.info(f"🌐 Website targeted: {url_input.strip()}")
    
    if st.button("🚀 Run Compliance Audit"):
        if not backend_engine:
            st.error("❌ Core analysis engine (`orchestrator.py`) could not be located in the folder structure.")
        else:
            with st.spinner("Extracting content and running live AI brand audit directly in the cloud..."):
                temp_image_path = None
                try:
                    text_to_check = ""
                    is_image = False
                    file_bytes = None
                    
                    if input_mode == "Live Website URL":
                        target_url = url_input.strip()
                        req = Request(target_url, headers={'User-Agent': 'Mozilla/5.0'})
                        with urlopen(req, timeout=10) as response:
                            raw_html = response.read().decode('utf-8', errors='ignore')
                        text_to_check = extract_text_from_html(raw_html)
                    else:
                        file_bytes = uploaded_file.getvalue()
                        
                        # Guardrail File Size Cap (15MB max for infrastructure safety)
                        if len(file_bytes) > 15 * 1024 * 1024:
                            st.error("❌ File size exceeds the 15MB system processing limit for live web builds.")
                            st.stop()
                            
                        file_ext = uploaded_file.name.lower().split('.')[-1]
                        
                        # Asset format pipeline routing
                        if file_ext == 'pdf':
                            text_to_check = extract_text_from_pdf(file_bytes)
                        elif file_ext in ['docx', 'doc']:
                            text_to_check = extract_text_from_docx(file_bytes)
                        elif file_ext in ['pptx', 'ppt']:
                            text_to_check = extract_text_from_pptx(file_bytes)
                        elif file_ext in ['png', 'jpg', 'jpeg']:
                            is_image = True
                            # Securely isolate uploaded image bytes on the cloud server using UUID names
                            temp_image_path = f"temp_{uuid.uuid4().hex}.{file_ext}"
                            with open(temp_image_path, "wb") as f:
                                f.write(file_bytes)
                            text_to_check = f"[Visual Layout Scan: {uploaded_file.name}]"
                        else:
                            try:
                                text_to_check = file_bytes.decode("utf-8", errors="ignore")
                            except Exception:
                                text_to_check = ""

                    if not is_image:
                        text_to_check = re.sub(r'\s+', ' ', text_to_check).strip()

                    if not is_image and not text_to_check.strip():
                        st.warning("⚠️ Could not extract any readable text copy out of the provided target.")
                    else:
                        # CALLS THE ORCHESTRATOR FUNCTION DIRECTLY IN MEMORY (No localhost connection required!)
                        result_data = backend_engine.run_brand_check(
                            text_to_check=text_to_check[:6000],
                            rules_context=rules_input,
                            image_path=temp_image_path
                        )
                        
                        # --- INTERFACE OUTPUT RENDERER ---
                        st.balloons()
                        
                        st.subheader("🔍 Contextual Verification Canvas")
                        issues = result_data.get("issues", [])
                        
                        if is_image:
                            st.image(uploaded_file, caption="Uploaded Visual Media Asset Canvas", use_container_width=True)
                        else:
                            marked_up_text = highlight_violations(text_to_check[:1200], issues)
                            st.markdown(f'<div style="background-color: #1E1E1E; padding: 20px; border-radius: 8px; border: 1px solid #333; line-height: 1.8; font-family: monospace; color: #E0E0E0;">{marked_up_text}...</div>', unsafe_allow_html=True)
                        st.markdown("---")
                        
                        st.subheader("📊 Executive Audit Summary")
                        st.write(result_data.get("summary", "No summary text generated."))
                        
                        st.subheader("⚠️ Detected Compliance Issues")
                        if not issues:
                            st.success("Perfect Compliance! No brand violations found matching your rules matrix.")
                        else:
                            for idx, issue in enumerate(issues, 1):
                                with st.expander(f"{idx}. {issue.get('issue_name', 'Style Deviation')}"):
                                    st.error(f"**Offending Element:** {issue.get('offending_text', 'N/A')}")
                                    st.success(f"**Required Correction:** {issue.get('corrected_text', 'N/A')}")
                                    st.caption(f"**Rule Violated:** {issue.get('brand_rule_violated', 'N/A')}")
                                    
                except Exception as e:
                    st.error(f"❌ An error occurred during analysis: {str(e)}")
                finally:
                    # Clean up temporary disk files instantly so cloud servers stay secure
                    if temp_image_path and os.path.exists(temp_image_path):
                        try:
                            os.remove(temp_image_path)
                        except Exception:
                            pass
