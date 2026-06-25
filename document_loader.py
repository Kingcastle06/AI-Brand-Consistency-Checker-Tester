import os
import pdfplumber
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_path: str) -> dict:
    """Extracts raw text and identifies unique fonts used in a PDF file."""
    text_content = []
    detected_fonts = set()

    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            # Pull text content cleanly
            page_text = page.extract_text()
            if page_text:
                text_content.append(page_text)

            # Look through page characters to find metadata on font names
            if page.chars:
                for char in page.chars:
                    font_name = char.get("fontname")
                    if font_name:
                        # Clean up sub-fonts or styles (e.g., "Arial-BoldMT" -> "Arial")
                        clean_font = font_name.split("-")[0].split(",")[0]
                        detected_fonts.add(clean_font)

    return {
        "text": "\n".join(text_content),
        "fonts": list(detected_fonts)
    }


def extract_text_from_docx(file_path: str) -> dict:
    """Extracts raw text and tracks font choices inside paragraphs and runs of a Word Doc."""
    doc = Document(file_path)
    text_content = []
    detected_fonts = set()

    for paragraph in doc.paragraphs:
        if paragraph.text:
            text_content.append(paragraph.text)
        
        # Look through specific text chunks to find inline font families
        for run in paragraph.runs:
            if run.font.name:
                detected_fonts.add(run.font.name)

    return {
        "text": "\n".join(text_content),
        "fonts": list(detected_fonts)
    }


def extract_text_from_pptx(file_path: str) -> dict:
    """Extracts text from all slide shapes and logs typography choices in a PowerPoint file."""
    prs = Presentation(file_path)
    text_content = []
    detected_fonts = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        text_content.append(paragraph.text)
                    for run in paragraph.runs:
                        if run.font.name:
                            detected_fonts.add(run.font.name)

    return {
        "text": "\n".join(text_content),
        "fonts": list(detected_fonts)
    }


def load_media_asset(file_path: str) -> dict:
    """Master entry point. Detects extension and maps to the correct parsing system."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Asset target not found at: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported document format extension: {ext}")